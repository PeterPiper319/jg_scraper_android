import os
import logging
from typing import Dict, Any

logger = logging.getLogger("MinerUParser")

class MinerUParser:
    @staticmethod
    def parse_document(file_path: str) -> Dict[str, Any]:
        """
        Parses scanned documents, DOCX, PPTX, XLSX using MinerU (magic-pdf)
        with robust pure-Python fallbacks for Office formats.
        """
        ext = os.path.splitext(file_path)[1].lower()
        logger.info(f"Parsing '{file_path}' using MinerU (format: {ext})...")
        
        try:
            # Check if magic_pdf is installed/available
            has_magic_pdf = False
            try:
                import magic_pdf
                has_magic_pdf = True
            except ImportError:
                pass

            if has_magic_pdf and ext == ".pdf":
                import tempfile
                import subprocess
                import shutil
                
                logger.info(f"magic_pdf is available. Running magic-pdf CLI on {file_path}")
                with tempfile.TemporaryDirectory() as temp_dir:
                    cmd = ["magic-pdf", "-p", file_path, "-o", temp_dir, "-m", "json"]
                    res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                    if res.returncode == 0:
                        base_name = os.path.splitext(os.path.basename(file_path))[0]
                        output_folder = os.path.join(temp_dir, base_name)
                        if os.path.exists(output_folder):
                            md_path = os.path.join(output_folder, f"{base_name}.md")
                            markdown_content = ""
                            if os.path.exists(md_path):
                                with open(md_path, "r", encoding="utf-8", errors="ignore") as f:
                                    markdown_content = f.read()
                            
                            # Read any generated table HTML files
                            tables_html = []
                            for root_dir, _, files in os.walk(output_folder):
                                for f in files:
                                    if f.endswith(".html"):
                                        try:
                                            with open(os.path.join(root_dir, f), "r", encoding="utf-8", errors="ignore") as tf:
                                                tables_html.append(tf.read())
                                        except Exception:
                                            pass
                            
                            return {
                                "markdown": markdown_content,
                                "tables_html": tables_html,
                                "metadata": {"parsed_via": "mineru_native"}
                            }
                    logger.warning(f"magic-pdf CLI failed (code {res.returncode}): {res.stderr}. Falling back to default parser.")

            # Fallback implementations based on file extension
            if ext == ".docx":
                return MinerUParser._parse_docx(file_path)
            elif ext in (".xlsx", ".xls"):
                return MinerUParser._parse_xlsx(file_path)
            elif ext == ".pptx":
                return MinerUParser._parse_pptx(file_path)
            else:
                logger.warning(f"Unsupported extension for MinerU: {ext}")
                return {"markdown": "", "tables_html": [], "metadata": {}}
        except Exception as e:
            logger.error(f"MinerU parsing failed for '{file_path}': {e}")
            return {"markdown": "", "tables_html": [], "metadata": {}}

    @staticmethod
    def _parse_docx(file_path: str) -> Dict[str, Any]:
        import docx
        doc = docx.Document(file_path)
        
        markdown_lines = []
        for p in doc.paragraphs:
            if p.text.strip():
                # Simple heading mapping
                if p.style.name.startswith("Heading"):
                    level = p.style.name.split()[-1]
                    hashes = "#" * (int(level) if level.isdigit() else 1)
                    markdown_lines.append(f"\n{hashes} {p.text.strip()}\n")
                else:
                    markdown_lines.append(p.text.strip())
                    
        tables_html = []
        for table in doc.tables:
            html = ["<table>"]
            for row in table.rows:
                html.append("  <tr>")
                for cell in row.cells:
                    # Clean the cell text
                    cell_text = cell.text.replace("\n", "<br>").strip()
                    html.append(f"    <td>{cell_text}</td>")
                html.append("  </tr>")
            html.append("</table>")
            tables_html.append("\n".join(html))
            
        return {
            "markdown": "\n".join(markdown_lines),
            "tables_html": tables_html,
            "metadata": {"parsed_via": "docx_fallback"}
        }

    @staticmethod
    def _parse_xlsx(file_path: str) -> Dict[str, Any]:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        
        markdown_lines = []
        tables_html = []
        
        for name in wb.sheetnames:
            sheet = wb[name]
            markdown_lines.append(f"\n## Sheet: {name}\n")
            
            html = [f"<h3>Sheet: {name}</h3>", "<table>"]
            sheet_text = []
            
            for row in sheet.iter_rows(values_only=True):
                if not any(row):  # Skip empty rows
                    continue
                html.append("  <tr>")
                row_vals = []
                for cell in row:
                    val = str(cell).strip() if cell is not None else ""
                    html.append(f"    <td>{val}</td>")
                    row_vals.append(val)
                html.append("  </tr>")
                sheet_text.append(" | ".join(row_vals))
                
            html.append("</table>")
            tables_html.append("\n".join(html))
            markdown_lines.append("\n".join(sheet_text))
            
        return {
            "markdown": "\n".join(markdown_lines),
            "tables_html": tables_html,
            "metadata": {"parsed_via": "openpyxl_fallback"}
        }

    @staticmethod
    def _parse_pptx(file_path: str) -> Dict[str, Any]:
        import pptx
        prs = pptx.Presentation(file_path)
        
        markdown_lines = []
        for idx, slide in enumerate(prs.slides):
            markdown_lines.append(f"\n## Slide {idx + 1}\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    markdown_lines.append(shape.text.strip())
                    
        return {
            "markdown": "\n".join(markdown_lines),
            "tables_html": [],
            "metadata": {"parsed_via": "python-pptx_fallback"}
        }
